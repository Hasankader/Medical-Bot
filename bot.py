import telebot
import json
import os
import re
import requests
from bs4 import BeautifulSoup
from deep_translator import GoogleTranslator
import pypdf
from pptx import Presentation
from PIL import Image
import io
import pytesseract
import os
# (باقي المكتبات مالتك تبقى مثل ما هي: telebot, json, re, إلخ...)

# هنا السيرفر راح يسحب التوكن بشكل سري بدون ما ينفضح بالكود
TOKEN = os.environ.get('BOT_TOKEN')
ADMIN_ID = 542533544  # الـ ID مالتك عادي يبقى هنا لأن ما يشكل خطر

bot = telebot.TeleBot(TOKEN)

# --- إدارة القاموس الطبي ---
def load_main_dict():
    if os.path.exists('medical_terms.json'):
        with open('medical_terms.json', 'r', encoding='utf-8') as f:
            try: return json.load(f)
            except: return {}
    return {}

def save_main_dict(data):
    with open('medical_terms.json', 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=4)

def fetch_reverso(word):
    url = f"https://context.reverso.net/translation/english-arabic/{word}"
    headers = {'User-Agent': 'Mozilla/5.0'}
    try:
        res = requests.get(url, headers=headers, timeout=2)
        if res.status_code == 200:
            soup = BeautifulSoup(res.content, 'html.parser')
            trans = soup.find('a', {'class': 'translation'})
            return trans.text.strip() if trans else None
    except: return None
    return None

def get_final_meaning(word, translator):
    personal = load_main_dict()
    w = word.lower()
    if w in personal: return personal[w]
    reverso_res = fetch_reverso(w)
    if reverso_res: return reverso_res
    try: return translator.translate(w)
    except: return w

# --- نظام الموافقة ---
@bot.message_handler(func=lambda m: '=' in m.text)
def handle_suggestion(message):
    try:
        parts = message.text.split('=')
        word = parts[0].strip().lower()
        suggested_trans = parts[1].strip()
        
        if message.chat.id == ADMIN_ID:
            data = load_main_dict()
            data[word] = suggested_trans
            save_main_dict(data)
            bot.reply_to(message, f"✅ تم تحديث القاموس فوراً:\n{word} = {suggested_trans}")
        else:
            reverso_help = fetch_reverso(word) or "لا توجد ترجمة واضحة"
            markup = telebot.types.InlineKeyboardMarkup()
            markup.add(
                telebot.types.InlineKeyboardButton("✅ موافق", callback_data=f"ok_{word}_{suggested_trans}"),
                telebot.types.InlineKeyboardButton("❌ رفض", callback_data=f"no_{word}")
            )
            bot.send_message(ADMIN_ID, f"🛡️ **طلب تصحيح:**\n`{word}` = `{suggested_trans}`\nرأي Reverso: {reverso_help}", reply_markup=markup)
            bot.reply_to(message, "تم إرسال الاقتراح للمشرف. ⏳")
    except: pass

@bot.callback_query_handler(func=lambda call: True)
def callback_query(call):
    if call.from_user.id != ADMIN_ID: return
    data_parts = call.data.split('_')
    action, word = data_parts[0], data_parts[1]
    if action == "ok":
        trans = data_parts[2]
        all_data = load_main_dict()
        all_data[word] = trans
        save_main_dict(all_data)
        bot.edit_message_text(f"✅ تم اعتماد: {word} = {trans}", call.message.chat.id, call.message.message_id)
    else:
        bot.edit_message_text(f"❌ تم رفض: {word}", call.message.chat.id, call.message.message_id)

# --- فلتر الترتيب الذكي للصور التشريحية ---
def organize_medical_text(raw_text):
    lines = raw_text.split('\n')
    main_paragraph = []
    labels = []
    
    for line in lines:
        # تنظيف الشخابيط اللي يسحبها الـ OCR
        line = re.sub(r'[^a-zA-Z0-9\s.,-]', '', line).strip()
        if len(line) < 3: continue
        
        # إذا السطر طويل (أكثر من 5 كلمات) نعتبره شرح أساسي
        if len(line.split()) > 5:
            main_paragraph.append(line)
        else:
            # إذا قصير، نعتبره تأشير (Label)
            labels.append(line)
            
    final_text = ""
    if main_paragraph:
        final_text += " ".join(main_paragraph) + ".\n\n"
    if labels:
        final_text += "Anatomical Labels:\n" + "\n".join(labels)
        
    return final_text

# --- معالجة وترجمة النصوص وبناء HTML ---
def process_full_lecture(slides_text, message):
    bot.reply_to(message, "جاري المعالجة والترجمة بشكل مرتب.. ⏳")
    try:
        translator = GoogleTranslator(source='auto', target='ar')
        pages_data = []
        all_words = list(set(re.findall(r'\b[a-zA-Z]{3,}\b', " ".join(slides_text).lower())))
        dynamic_dict = {w: get_final_meaning(w, translator) for w in all_words}
        
        for slide in slides_text:
            sentences = re.split(r'(?<=[.!?])\s+|\n+', slide)
            slide_content = []
            for s in sentences:
                if len(s.strip()) > 2:
                    try: trans = translator.translate(s)
                    except: trans = s
                    slide_content.append({"text": s, "trans": trans})
            if slide_content: pages_data.append({"content": slide_content})

        with open('template.html', 'r', encoding='utf-8') as f: template = f.read()
        final_html = template.replace('{{PAGES_DATA}}', json.dumps(pages_data, ensure_ascii=False)).replace('{{DICTIONARY_DATA}}', json.dumps(dynamic_dict, ensure_ascii=False))

        out_name = f"Lecture_{message.chat.id}.html"
        with open(out_name, 'w', encoding='utf-8') as f: f.write(final_html)
        bot.send_document(message.chat.id, open(out_name, 'rb'), caption="تمت الترجمة بترتيب دقيق! ✅")
        os.remove(out_name)
    except Exception as e: bot.reply_to(message, f"خطأ في الترجمة: {e}")

# --- استقبال الملفات والصور ---
@bot.message_handler(content_types=['document', 'photo'])
def handle_file(message):
    try:
        bot.reply_to(message, "استلمت الصورة.. جاري قراءة التأشيرات بالترتيب ⏳")
        if message.content_type == 'photo':
            file_id = message.photo[-1].file_id
            ext = 'jpg'
        else:
            file_id = message.document.file_id
            ext = message.document.file_name.split('.')[-1].lower()

        file_info = bot.get_file(file_id)
        path = f"temp_{message.chat.id}.{ext}"
        with open(path, 'wb') as f: f.write(bot.download_file(file_info.file_path))
        
        slides = []
        
        # استخدام psm 4 للقراءة سطر بسطر (من الأعلى للأسفل)
        if ext in ['jpg', 'jpeg', 'png']:
            img = Image.open(path)
            raw_text = pytesseract.image_to_string(img, lang='eng', config='--psm 4')
            organized = organize_medical_text(raw_text)
            if organized.strip(): slides.append(organized)
            
        elif ext == 'pdf':
            reader = pypdf.PdfReader(path)
            for page in reader.pages:
                t = page.extract_text() or ""
                try: 
                    for img_obj in page.images:
                        img = Image.open(io.BytesIO(img_obj.data))
                        raw = pytesseract.image_to_string(img, lang='eng', config='--psm 4')
                        t += "\n" + organize_medical_text(raw)
                except: pass
                if t.strip(): slides.append(t.strip())
                
        elif ext == 'pptx':
            prs = Presentation(path)
            for s in prs.slides:
                t = "\n".join([sh.text for sh in s.shapes if hasattr(sh, "text")])
                for sh in s.shapes:
                    if sh.shape_type == 13: 
                        try: 
                            raw = pytesseract.image_to_string(Image.open(io.BytesIO(sh.image.blob)), lang='eng', config='--psm 4')
                            t += "\n" + organize_medical_text(raw)
                        except: pass
                if t.strip(): slides.append(t.strip())
        
        os.remove(path)
        if slides: process_full_lecture(slides, message)
        else: bot.reply_to(message, "❌ لم أتمكن من قراءة أي نص. حاول تقريب الصورة.")
    except Exception as e: bot.reply_to(message, f"خطأ بالملف: {e}")

@bot.message_handler(commands=['start'])
def start(message):
    bot.reply_to(message, "أهلاً بك! ارسل ملزمة أو صورة تشريحية ليتم ترتيبها وترجمتها.")

bot.polling(none_stop=True)
